#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
عرض تقديمي: التفاوض بين Microsoft و Nokia
Presentation: Microsoft-Nokia Negotiation Case Study
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """إنشاء عرض تقديمي PowerPoint عن التفاوض بين Microsoft و Nokia"""
    
    # إنشاء عرض تقديمي جديد
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # الألوان المستخدمة
    MICROSOFT_BLUE = RGBColor(0, 120, 215)
    NOKIA_BLUE = RGBColor(0, 92, 171)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(128, 128, 128)
    GREEN = RGBColor(16, 185, 129)
    RED = RGBColor(239, 68, 68)
    
    def add_title_slide(title, subtitle):
        """إضافة شريحة العنوان"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # خلفية زرقاء متدرجة
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = MICROSOFT_BLUE
        
        # العنوان الرئيسي
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # العنوان الفرعي
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # التاريخ
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = "نوفمبر 2025"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = WHITE
        date_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_items, bg_color=WHITE):
        """إضافة شريحة محتوى"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # الخلفية
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # العنوان
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = MICROSOFT_BLUE
        title_para.alignment = PP_ALIGN.RIGHT
        
        # خط فاصل
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(0.5), Inches(1.4), Inches(9), Inches(0)
        )
        line.line.color.rgb = MICROSOFT_BLUE
        line.line.width = Pt(3)
        
        # المحتوى
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = BLACK
            p.alignment = PP_ALIGN.RIGHT
            p.level = 0
            p.space_before = Pt(12)
    
    # شريحة 1: الغلاف
    add_title_slide(
        "التفاوض بين Microsoft و Nokia",
        "دراسة حالة في الاستحواذ الاستراتيجي"
    )
    
    # شريحة 2: المشكلة
    add_content_slide(
        "المشكلة: تحديات السوق",
        [
            "🔴 Microsoft: فشل في دخول سوق الهواتف الذكية بقوة",
            "",
            "🔴 Nokia: انهيار مبيعات الهواتف بسبب منافسة Android و iOS",
            "",
            "💡 السؤال: كيف يمكن للشركتين التعاون لمواجهة التحديات؟",
            "",
            "📊 الحاجة المتبادلة:",
            "   • Microsoft تحتاج خبرة تصنيع الأجهزة",
            "   • Nokia تحتاج نظام تشغيل قوي ودعم مالي"
        ]
    )
    
    # شريحة 3: الخلفية التاريخية
    add_content_slide(
        "الخلفية التاريخية",
        [
            "📅 بداية 2013:",
            "   • Nokia تواجه أزمة كبيرة في مبيعات الهواتف",
            "   • انهيار حصتها السوقية بعد ظهور iPhone و Android",
            "",
            "📱 Microsoft:",
            "   • محاولات فاشلة لدخول سوق الهواتف بـ Windows Phone",
            "   • حصة سوقية ضعيفة جداً (أقل من 3%)",
            "",
            "🤝 بداية المفاوضات:",
            "   • Microsoft تسعى للاستحواذ على قسم الأجهزة في Nokia",
            "   • الهدف: دمج الأجهزة مع نظام Windows Phone"
        ]
    )
    
    # شريحة 4: تحليل Microsoft
    add_content_slide(
        "الطرف الأول: Microsoft",
        [
            "🎯 الهدف الاستراتيجي:",
            "   • دخول سوق الهواتف عبر شراء علامة تجارية قوية",
            "   • تجنب البدء من الصفر",
            "",
            "💼 تكتيكات التفاوض:",
            "",
            "1️⃣ المشاركة:",
            "   • عرض بقاء إدارة Nokia مؤقتاً",
            "   • ضمان استمرارية العمليات",
            "",
            "2️⃣ المكاشفة:",
            "   • وضوح في الأرقام المالية والأرباح المتوقعة",
            "",
            "3️⃣ الضغط:",
            "   • التلويح ببدائل أخرى للضغط على Nokia"
        ]
    )
    
    # شريحة 5: نتيجة Microsoft
    add_content_slide(
        "نتيجة Microsoft",
        [
            "✅ النجاح المالي:",
            "   • إتمام الصفقة مقابل 7.2 مليار دولار (2014)",
            "   • الاستحواذ على قسم الأجهزة والخدمات",
            "   • الحصول على براءات اختراع Nokia",
            "",
            "❌ الفشل السوقي:",
            "   • Windows Phone فشل في منافسة Android و iOS",
            "   • انخفاض مستمر في الحصة السوقية",
            "   • إغلاق قسم الهواتف في 2016",
            "",
            "💰 الخسارة:",
            "   • شطب 7.6 مليار دولار في 2015",
            "   • تسريح آلاف الموظفين"
        ]
    )
    
    # شريحة 6: تحليل Nokia
    add_content_slide(
        "الطرف الثاني: Nokia",
        [
            "🎯 الهدف الاستراتيجي:",
            "   • إنقاذ الشركة من الإفلاس",
            "   • الحفاظ على الاسم التجاري والموظفين",
            "",
            "💼 تكتيكات التفاوض:",
            "",
            "1️⃣ كسب الوقت:",
            "   • تأخير الموافقة لجمع عروض بديلة",
            "   • زيادة القيمة التفاوضية",
            "",
            "2️⃣ التراجع المؤقت:",
            "   • بيع قسم الهواتف فقط",
            "   • الاحتفاظ بأقسام الشبكات والبحث والتطوير",
            "",
            "3️⃣ الجانب الإنساني:",
            "   • التركيز على إنقاذ آلاف الوظائف"
        ]
    )
    
    # شريحة 7: نتيجة Nokia
    add_content_slide(
        "نتيجة Nokia",
        [
            "✅ النجاح المالي:",
            "   • الحصول على 7.2 مليار دولار نقداً",
            "   • سداد الديون وإعادة الهيكلة",
            "",
            "✅ التحول الاستراتيجي:",
            "   • التركيز على Nokia Networks (معدات الشبكات)",
            "   • نمو قوي في مجال البنية التحتية للاتصالات",
            "   • الاستثمار في تقنية 5G",
            "",
            "📈 النتيجة الحالية:",
            "   • Nokia أصبحت من أكبر موردي معدات الشبكات عالمياً",
            "   • عودة قوية للسوق في مجال جديد",
            "   • استقرار مالي وربحية مستدامة"
        ]
    )
    
    # شريحة 8: الدروس المستفادة
    add_content_slide(
        "الدروس التفاوضية المستفادة",
        [
            "💡 الدرس الأول: أهمية الشفافية",
            "   • الوضوح في الأهداف والأرقام يسرع التفاوض",
            "",
            "💡 الدرس الثاني: التخطيط الاستراتيجي",
            "   • فهم نقاط القوة والضعف لكلا الطرفين",
            "",
            "💡 الدرس الثالث: المرونة",
            "   • Nokia نجحت بالتركيز على مجال جديد",
            "   • عدم التمسك بالماضي",
            "",
            "💡 الدرس الرابع: دور الضغط والتعاطف",
            "   • استخدام عوامل متعددة في التفاوض",
            "",
            "⚠️ الدرس الخامس: النجاح المالي ≠ النجاح السوقي",
            "   • الصفقة الجيدة تحتاج تنفيذ استراتيجي ناجح"
        ]
    )
    
    # شريحة 9: المقارنة
    add_content_slide(
        "مقارنة النتائج",
        [
            "📊 Microsoft:",
            "   ✅ استحواذ ناجح من الناحية القانونية",
            "   ❌ فشل في تحقيق الأهداف السوقية",
            "   ❌ خسارة مالية كبيرة (7.6 مليار دولار)",
            "   📉 خروج من سوق الهواتف",
            "",
            "📊 Nokia:",
            "   ✅ إنقاذ الشركة من الإفلاس",
            "   ✅ تحول استراتيجي ناجح",
            "   ✅ نمو قوي في مجال الشبكات",
            "   📈 عودة قوية للربحية",
            "",
            "🏆 الفائز الحقيقي: Nokia (على المدى الطويل)"
        ]
    )
    
    # شريحة 10: الخاتمة
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NOKIA_BLUE
    
    # العنوان
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "الخاتمة"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # المحتوى
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(2))
    text_frame = content_box.text_frame
    text_frame.text = "صفقة ناجحة مالياً لكلا الطرفين\nلكنها تذكرنا بأهمية التكيف مع السوق\nوالتنفيذ الاستراتيجي الناجح"
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(12)
    
    # شكراً
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "شكراً لكم"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(36)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = WHITE
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # حفظ العرض التقديمي
    output_file = "/vercel/sandbox/Microsoft_Nokia_Negotiation.pptx"
    prs.save(output_file)
    print(f"✅ تم إنشاء العرض التقديمي بنجاح: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
